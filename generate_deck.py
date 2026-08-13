import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# CAPSULE TRACKER - 14-SLIDE PREMIUM POWERPOINT GENERATOR
# ==============================================================================

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 Widescreen
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# --- BRAND COLOR PALETTE ---
BG_DARK      = RGBColor(15, 23, 42)     # #0F172A (Deep Slate)
CARD_BG      = RGBColor(29, 40, 60)     # #1D283C (Card Dark)
CARD_BORDER  = RGBColor(51, 65, 85)    # #334155
TEXT_WHITE   = RGBColor(248, 250, 252) # #F8FAFC
TEXT_MUTED   = RGBColor(148, 163, 184) # #94A3B8
ACCENT_CYAN  = RGBColor(56, 189, 248)  # #38BDF8 (Primary Accent)
ACCENT_GREEN = RGBColor(34, 197, 94)   # #22C55E (Safe / Live)
ACCENT_RED   = RGBColor(239, 68, 68)    # #EF4444 (SOS / Alert)

def add_blank_slide():
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return slide

def add_header(slide, title_text, category_text):
    # Category Tag (13pt)
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(13)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ACCENT_CYAN
    p_cat.font.name = "Inter"

    # Slide Title (28pt)
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.5), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.font.name = "Inter"

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
    return card

def add_image_container(slide, left, top, width, height, label="IMAGE / SCREENSHOT PLACEHOLDER", image_path=None):
    """Inserts an image if available; otherwise builds a styled image placeholder card."""
    if image_path and os.path.exists(image_path):
        return slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
    
    # Outer Placeholder Box
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(20, 30, 48)
    frame.line.color.rgb = ACCENT_CYAN
    frame.line.width = Pt(1.5)
    
    # Center Label Text
    tf = frame.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"🖼️ {label}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.alignment = PP_ALIGN.CENTER
    return frame

# ==============================================================================
# SLIDE 1: TITLE SLIDE
# ==============================================================================
s1 = add_blank_slide()

bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(0.18), Inches(3.6))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_CYAN
bar.line.fill.background()

t_box = s1.shapes.add_textbox(Inches(1.4), Inches(1.8), Inches(10.5), Inches(3.8))
tf = t_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "CAPSULE TRACKER"
p1.font.size = Pt(46)
p1.font.bold = True
p1.font.color.rgb = TEXT_WHITE
p1.font.name = "Inter"

p2 = tf.add_paragraph()
p2.text = "Child Safety & Off-Grid Telemetry System"
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = ACCENT_CYAN
p2.font.name = "Inter"
p2.space_before = Pt(10)

p3 = tf.add_paragraph()
p3.text = "Real-Time GPS Tracking, Biometric Vitals Monitoring & Instant Tamper Alerts"
p3.font.size = Pt(16)
p3.font.color.rgb = TEXT_MUTED
p3.font.name = "Inter"
p3.space_before = Pt(18)

card_f = add_card(s1, 1.4, 5.5, 5.8, 0.85, bg_color=CARD_BG, border_color=ACCENT_GREEN)
tf_f = card_f.text_frame
p_f = tf_f.paragraphs[0]
p_f.text = "🟢 Status: Live Backend Ingestion & Dashboard Operational"
p_f.font.size = Pt(14)
p_f.font.color.rgb = ACCENT_GREEN
p_f.font.bold = True
p_f.alignment = PP_ALIGN.CENTER

# ==============================================================================
# SLIDE 2: PROBLEM STATEMENT & MOTIVATION
# ==============================================================================
s2 = add_blank_slide()
add_header(s2, "Problem Statement & Project Motivation", "Background & Challenge")

c2_1 = add_card(s2, 0.8, 1.7, 5.6, 5.2)
tf2_1 = c2_1.text_frame
tf2_1.word_wrap = True
tf2_1.margin_left = Inches(0.3)
tf2_1.margin_top = Inches(0.3)

p = tf2_1.paragraphs[0]
p.text = "⚠️ Existing Safety Limitations"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_RED

bullets2_1 = [
    "Cellular Dependence: Standard GPS trackers fail entirely in remote or off-grid areas without SIM signal.",
    "No Removal Detection: Traditional wearables provide no alert when unbuckled or forcibly cut.",
    "Lack of Vitals Context: Basic location trackers fail to capture stress signals like elevated heart rates."
]
for b in bullets2_1:
    p = tf2_1.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(14)

c2_2 = add_card(s2, 6.8, 1.7, 5.7, 5.2)
tf2_2 = c2_2.text_frame
tf2_2.word_wrap = True
tf2_2.margin_left = Inches(0.3)
tf2_2.margin_top = Inches(0.3)

p = tf2_2.paragraphs[0]
p.text = "💡 The Capsule Tracker Solution"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

bullets2_2 = [
    "Long-Range LoRa Mesh: Transmits telemetry parameters up to several kilometers without cellular networks.",
    "Optical Strap Tamper Sensor: Triggers instantaneous alerts if the wristband is opened.",
    "Integrated Biometrics: Monitors real-time Heart Rate (BPM) and Blood Oxygen (SpO2).",
    "Automated Geofence: Instantly flags out-of-bounds occurrences via backend algorithms."
]
for b in bullets2_2:
    p = tf2_2.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(14)

# ==============================================================================
# SLIDE 3: SYSTEM ARCHITECTURE OVERVIEW
# ==============================================================================
s3 = add_blank_slide()
add_header(s3, "High-Level System Architecture", "System Topology")

add_image_container(s3, 0.8, 1.7, 5.6, 5.2, "SYSTEM ARCHITECTURE DIAGRAM")

c3 = add_card(s3, 6.8, 1.7, 5.7, 5.2)
tf3 = c3.text_frame
tf3.word_wrap = True
tf3.margin_left = Inches(0.3)
tf3.margin_top = Inches(0.3)

p = tf3.paragraphs[0]
p.text = "🔄 End-to-End Data Pipeline"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

pipeline_steps = [
    "1. Smart Wristband Node: Captures GPS coordinates, heart rate, SpO2, SOS button state, and optical tamper status.",
    "2. LoRa Telemetry Transmission: Packets are broadcast over long-range radio frequencies to the central gateway.",
    "3. PHP REST Ingestion Engine: `receive_telemetry.php` receives JSON payloads, executes validation, and processes geofencing.",
    "4. MySQL Relational Storage: Ingested telemetry is committed to normalized SQL tables (`telemetry_logs` & `alert_logs`).",
    "5. Async Live Dashboard: AJAX long-polling updates Google Maps Dark Mode UI every 3 seconds."
]
for step in pipeline_steps:
    p = tf3.add_paragraph()
    p.text = step
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(10)

# ==============================================================================
# SLIDE 4: HARDWARE WEARABLE NODE DESIGN
# ==============================================================================
s4 = add_blank_slide()
add_header(s4, "Smart Wristband Hardware & Sensor Suite", "Wearable Layer")

c4 = add_card(s4, 0.8, 1.7, 5.6, 5.2)
tf4 = c4.text_frame
tf4.word_wrap = True
tf4.margin_left = Inches(0.3)
tf4.margin_top = Inches(0.3)

p = tf4.paragraphs[0]
p.text = "🎛️ On-Board Components"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

hw_specs = [
    "LoRa SX1278 Transceiver: 433/868MHz frequency band for long-range transmission.",
    "GPS Module (NEO-6M): High-precision coordinate positioning.",
    "PPG Pulse Sensor (MAX30102): Real-time optical heart rate and SpO2 calculation.",
    "Optical Strap Sensor: Detects physical circuit broken / removal status.",
    "Panic Button (SOS): Tactile switch triggering immediate high-priority override alerts."
]
for hw in hw_specs:
    p = tf4.add_paragraph()
    p.text = f"• {hw}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

add_image_container(s4, 6.8, 1.7, 5.7, 5.2, "WRISTBAND HARDWARE SCHEMATIC")

# ==============================================================================
# SLIDE 5: DATABASE ARCHITECTURE (MYSQL)
# ==============================================================================
s5 = add_blank_slide()
add_header(s5, "Relational Database Schema (MySQL)", "Data Storage & Ingestion")

c5_1 = add_card(s5, 0.8, 1.7, 5.6, 5.2)
tf5_1 = c5_1.text_frame
tf5_1.word_wrap = True
tf5_1.margin_left = Inches(0.3)
tf5_1.margin_top = Inches(0.3)

p = tf5_1.paragraphs[0]
p.text = "📋 `telemetry_logs` Schema"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

t_schema = [
    "id: INT AUTO_INCREMENT PRIMARY KEY",
    "device_id: VARCHAR(50) DEFAULT 'WRISTBAND_01'",
    "latitude & longitude: DECIMAL(10,8) / (11,8)",
    "speed: FLOAT (Movement tracking)",
    "battery_level: INT (0-100%)",
    "heart_rate & spo2: INT (Biometric values)",
    "tamper_status & sos_active: TINYINT(1)",
    "created_at: TIMESTAMP DEFAULT CURRENT_TIMESTAMP"
]
for t in t_schema:
    p = tf5_1.add_paragraph()
    p.text = f"• {t}"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

c5_2 = add_card(s5, 6.8, 1.7, 5.7, 5.2)
tf5_2 = c5_2.text_frame
tf5_2.word_wrap = True
tf5_2.margin_left = Inches(0.3)
tf5_2.margin_top = Inches(0.3)

p = tf5_2.paragraphs[0]
p.text = "🚨 `alert_logs` Schema"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_RED

a_schema = [
    "id: INT AUTO_INCREMENT PRIMARY KEY",
    "device_id: VARCHAR(50) NOT NULL",
    "alert_type: VARCHAR(50) ('GEOFENCE_BREACH', 'STRAP_TAMPER', 'SOS_EMERGENCY')",
    "message: TEXT (Detailed emergency description)",
    "latitude & longitude: DECIMAL(10,8) / (11,8)",
    "is_resolved: TINYINT(1) DEFAULT 0",
    "created_at: TIMESTAMP DEFAULT CURRENT_TIMESTAMP"
]
for a in a_schema:
    p = tf5_2.add_paragraph()
    p.text = f"• {a}"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(10)

# ==============================================================================
# SLIDE 6: GEOFENCING MATHEMATICS & ALGORITHM
# ==============================================================================
s6 = add_blank_slide()
add_header(s6, "Real-Time Geofencing Engine", "Algorithmic Logic")

c6_1 = add_card(s6, 0.8, 1.7, 5.6, 5.2)
tf6_1 = c6_1.text_frame
tf6_1.word_wrap = True
tf6_1.margin_left = Inches(0.3)
tf6_1.margin_top = Inches(0.3)

p = tf6_1.paragraphs[0]
p.text = "📐 Haversine Distance Formula"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

p2 = tf6_1.add_paragraph()
p2.text = "Calculates exact great-circle distance between two points on a sphere:"
p2.font.size = Pt(14)
p2.font.color.rgb = TEXT_MUTED
p2.space_before = Pt(8)

haversine_points = [
    "Δlat = lat2 - lat1 | Δlon = lon2 - lon1",
    "a = sin²(Δlat/2) + cos(lat1)·cos(lat2)·sin²(Δlon/2)",
    "c = 2 · atan2( √a, √(1−a) )",
    "Distance = R · c  (where R = 6,371,000 meters)"
]
for hp in haversine_points:
    p = tf6_1.add_paragraph()
    p.text = f"• {hp}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(10)

c6_2 = add_card(s6, 6.8, 1.7, 5.7, 5.2)
tf6_2 = c6_2.text_frame
tf6_2.word_wrap = True
tf6_2.margin_left = Inches(0.3)
tf6_2.margin_top = Inches(0.3)

p = tf6_2.paragraphs[0]
p.text = "🛡️ Safe Zone Enforcement"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

gf_rules = [
    "Dynamic Radius Boundary: Preset home/school center point with configurable safe radius (e.g., 500m).",
    "Automated Breach Detection: Evaluates incoming coordinates against safe boundary during telemetry parsing.",
    "Instant Alert Dispatch: Automatically logs 'GEOFENCE_BREACH' entry into `alert_logs` if computed distance exceeds radius."
]
for gfr in gf_rules:
    p = tf6_2.add_paragraph()
    p.text = f"• {gfr}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(14)

# ==============================================================================
# SLIDE 7: EMERGENCY & TAMPER DETECTION SYSTEM
# ==============================================================================
s7 = add_blank_slide()
add_header(s7, "Emergency & Tamper Alert Ingestion", "Security & Protection")

c7 = add_card(s7, 0.8, 1.7, 5.6, 5.2, border_color=ACCENT_RED)
tf7 = c7.text_frame
tf7.word_wrap = True
tf7.margin_left = Inches(0.3)
tf7.margin_top = Inches(0.3)

p = tf7.paragraphs[0]
p.text = "⚡ Real-Time Threat Handling"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_RED

threat_list = [
    "Strap Tamper Trigger: Instant detection when the optical circuit opens. `receive_telemetry.php` immediately logs 'STRAP_TAMPER'.",
    "Panic SOS Trigger: Manual activation by child creates high-priority 'SOS_EMERGENCY' alert.",
    "Visual UI Override: Flashes red emergency badge on the frontend monitoring system.",
    "Unresolved Alert Queue: Persists warning state until administrator clears the event."
]
for th in threat_list:
    p = tf7.add_paragraph()
    p.text = f"• {th}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

add_image_container(s7, 6.8, 1.7, 5.7, 5.2, "ALERT BADGE OVERRIDE SCREENSHOT")

# ==============================================================================
# SLIDE 8: REAL-TIME MONITORING DASHBOARD (UI)
# ==============================================================================
s8 = add_blank_slide()
add_header(s8, "Web Control Center & Interactive Map", "Frontend Dashboard")

add_image_container(s8, 0.8, 1.7, 5.6, 5.2, "LIVE DASHBOARD SCREENSHOT (INDEX.HTML)")

c8 = add_card(s8, 6.8, 1.7, 5.7, 5.2)
tf8 = c8.text_frame
tf8.word_wrap = True
tf8.margin_left = Inches(0.3)
tf8.margin_top = Inches(0.3)

p = tf8.paragraphs[0]
p.text = "🖥️ Dashboard Interface (`index.html`)"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

dash_features = [
    "Google Maps Dark Mode: High contrast, dark styled vector map centered on live coordinates.",
    "Async 3-Second Poller: `fetchDashboardData()` retrieves telemetry smooth without page reloads.",
    "Biometric Cards: Live Heart Rate (BPM), SpO2 (%), and Wearable Battery percentage.",
    "Quick Emergency Controls: Ping Wristband transmitter and Remote Alarm Siren trigger."
]
for df in dash_features:
    p = tf8.add_paragraph()
    p.text = f"• {df}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

# ==============================================================================
# SLIDE 9: WEARABLE TELEMETRY SIMULATOR
# ==============================================================================
s9 = add_blank_slide()
add_header(s9, "Hardware Telemetry Simulator Environment", "Testing & Emulation")

c9 = add_card(s9, 0.8, 1.7, 5.6, 5.2)
tf9 = c9.text_frame
tf9.word_wrap = True
tf9.margin_left = Inches(0.3)
tf9.margin_top = Inches(0.3)

p = tf9.paragraphs[0]
p.text = "🛠️ Telemetry Simulator (`simulator.html`)"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

sim_features = [
    "Coordinate Manipulation: Test custom GPS points to simulate geofence boundaries.",
    "Biometric Sliders: Dynamic adjustment of simulated Heart Rate (BPM) and SpO2 levels.",
    "Hardware State Toggles: Instantly force SOS Active or Strap Cut states.",
    "API Response Console: Displays real-time JSON responses returned by PHP scripts."
]
for sf in sim_features:
    p = tf9.add_paragraph()
    p.text = f"• {sf}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

add_image_container(s9, 6.8, 1.7, 5.7, 5.2, "SIMULATOR UI SCREENSHOT (SIMULATOR.HTML)")

# ==============================================================================
# SLIDE 10: BACKEND REST API ENDPOINTS
# ==============================================================================
s10 = add_blank_slide()
add_header(s10, "RESTful Backend API Infrastructure", "PHP API Architecture")

c10_1 = add_card(s10, 0.8, 1.7, 5.6, 5.2)
tf10_1 = c10_1.text_frame
tf10_1.word_wrap = True
tf10_1.margin_left = Inches(0.3)
tf10_1.margin_top = Inches(0.3)

p = tf10_1.paragraphs[0]
p.text = "📥 `receive_telemetry.php`"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

api1_bullets = [
    "Method: POST (JSON Payload Ingestion)",
    "Validates input payload fields and checks bounds.",
    "Executes Haversine Geofence comparison.",
    "Inserts packet into `telemetry_logs`.",
    "Generates emergency alerts in `alert_logs` if tamper or SOS active."
]
for b in api1_bullets:
    p = tf10_1.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(10)

c10_2 = add_card(s10, 6.8, 1.7, 5.7, 5.2)
tf10_2 = c10_2.text_frame
tf10_2.word_wrap = True
tf10_2.margin_left = Inches(0.3)
tf10_2.margin_top = Inches(0.3)

p = tf10_2.paragraphs[0]
p.text = "📤 `get_dashboard_data.php`"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

api2_bullets = [
    "Method: GET (JSON Response)",
    "Queries latest record from `telemetry_logs`.",
    "Fetches active unresolved warnings from `alert_logs`.",
    "Returns synchronized JSON telemetry state to dashboard JS engine."
]
for b in api2_bullets:
    p = tf10_2.add_paragraph()
    p.text = f"• {b}"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

# ==============================================================================
# SLIDE 11: KEY TECHNICAL CHALLENGES & RESOLUTIONS
# ==============================================================================
s11 = add_blank_slide()
add_header(s11, "Technical Challenges & Solutions Resolved", "Debugging & Optimization")

c11 = add_card(s11, 0.8, 1.7, 11.7, 5.2)
tf11 = c11.text_frame
tf11.word_wrap = True
tf11.margin_left = Inches(0.4)
tf11.margin_top = Inches(0.3)

p = tf11.paragraphs[0]
p.text = "🔧 Engineering Resolutions"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

fixes = [
    ("Database Column Mismatch (SQL State 42S22)", "Resolved missing 'device_id' column issue in `alert_logs` table via automated migration scripts."),
    ("Async Dynamic Map Polling", "Replaced static random telemetry simulation in JS with an active 3-second fetch engine connected to `get_dashboard_data.php`."),
    ("Google Maps Marker Repositioning", "Implemented smooth `panTo()` map re-centering without full iframe/DOM re-renders."),
    ("Visual Alert Badge Override", "Configured conditional CSS logic to switch status badges dynamically based on backend SOS/Tamper flags.")
]

for title, desc in fixes:
    p = tf11.add_paragraph()
    p.text = f"• {title}: {desc}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(14)

# ==============================================================================
# SLIDE 12: DEVELOPMENT ROADMAP & MILESTONES
# ==============================================================================
s12 = add_blank_slide()
add_header(s12, "Development Roadmap & Future Expansion", "Project Phasing")

c12_1 = add_card(s12, 0.8, 1.7, 5.6, 5.2, border_color=ACCENT_GREEN)
tf12_1 = c12_1.text_frame
tf12_1.word_wrap = True
tf12_1.margin_left = Inches(0.3)
tf12_1.margin_top = Inches(0.3)

p = tf12_1.paragraphs[0]
p.text = "✅ Phase 1: Completed"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_GREEN

phase1 = [
    "MySQL database schema & connection setup",
    "Haversine geofencing & backend telemetry ingestion API",
    "Interactive telemetry simulator interface",
    "Live dark-mode dashboard with async polling",
    "Real-time visual alert system"
]
for p1_item in phase1:
    p = tf12_1.add_paragraph()
    p.text = f"✓ {p1_item}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

c12_2 = add_card(s12, 6.8, 1.7, 5.7, 5.2, border_color=ACCENT_CYAN)
tf12_2 = c12_2.text_frame
tf12_2.word_wrap = True
tf12_2.margin_left = Inches(0.3)
tf12_2.margin_top = Inches(0.3)

p = tf12_2.paragraphs[0]
p.text = "🚀 Phase 2: Next Steps"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

phase2 = [
    "Geofence UI Configuration (`Geofencing.html`)",
    "Historical Breadcrumb Route Playback (`Tracking_History.html`)",
    "Device & Threshold Settings UI (`Setting.html`)",
    "Physical LoRa SX1278 ESP32 Hardware Deployment"
]
for p2_item in phase2:
    p = tf12_2.add_paragraph()
    p.text = f"→ {p2_item}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(14)

# ==============================================================================
# SLIDE 13: LIVE DEMONSTRATION WORKFLOW
# ==============================================================================
s13 = add_blank_slide()
add_header(s13, "Live System Demonstration Workflow", "Execution Plan")

c13 = add_card(s13, 0.8, 1.7, 11.7, 5.2)
tf13 = c13.text_frame
tf13.word_wrap = True
tf13.margin_left = Inches(0.4)
tf13.margin_top = Inches(0.3)

p = tf13.paragraphs[0]
p.text = "🧪 Live Test Procedure"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN

demo_steps = [
    "Step 1: Open Live Dashboard (`index.html`) in Primary Window.",
    "Step 2: Launch Telemetry Simulator (`simulator.html`) in Secondary Window.",
    "Step 3: Modify Latitude/Longitude to push device outside safe radius and observe instant Geofence Alert.",
    "Step 4: Toggle Strap Cut / SOS Emergency state in Simulator.",
    "Step 5: Verify automatic MySQL database commit and dynamic UI status badge shift within 3 seconds."
]
for ds in demo_steps:
    p = tf13.add_paragraph()
    p.text = ds
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(14)

# ==============================================================================
# SLIDE 14: CONCLUSION & Q&A
# ==============================================================================
s14 = add_blank_slide()

c14 = add_card(s14, 1.5, 1.5, 10.3, 4.5, bg_color=CARD_BG, border_color=ACCENT_CYAN)
tf14 = c14.text_frame
tf14.word_wrap = True
tf14.margin_left = Inches(0.5)
tf14.margin_top = Inches(0.5)

p = tf14.paragraphs[0]
p.text = "THANK YOU!"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = TEXT_WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf14.add_paragraph()
p2.text = "Capsule Tracker: Smart Off-Grid Child Safety & Telemetry"
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT_CYAN
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)

p3 = tf14.add_paragraph()
p3.text = "Questions & Demonstration"
p3.font.size = Pt(18)
p3.font.color.rgb = TEXT_MUTED
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(24)

# --- SAVE PRESENTATION ---
output_filename = "Capsule_Tracker_Presentation_V2.pptx"
prs.save(output_filename)
print(f"✨ Presentation successfully generated and saved as '{output_filename}'!")